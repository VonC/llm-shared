"""Slides 9-16 of the llm-shared deck.

They cover the workflow overview, the two phases, groundhog, commits and
release, security, benefits and the conclusion.

Every pixel coordinate is read from the matching slide in
docs/llm-shared_presentation.html; the Theme maps 1280x720 px onto a
13.333in x 7.5in 16:9 slide, so the numbers transfer without re-layout.
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from tools.html_to_pptx.deck_parts import (
    SecCellStyle as Sec,
)
from tools.html_to_pptx.deck_parts import (
    StepRow,
    brand_top_left,
    code_run,
    sec_cell,
    step_row,
    wf_row,
)
from tools.html_to_pptx.deck_theme import (
    BLUE,
    BRAND,
    BRAND_SUB,
    CODE,
    DARK,
    GRAY,
    GREEN,
    LIGHT,
    LIGHTBLUE,
    MUTED,
    ORANGE,
    TEAL,
    TEXT,
    THEME,
    TMUTED,
    TSOFT,
    WARM,
    WHITE,
)
from tools.html_to_pptx.pptx_helpers import (
    Area,
    BoxStyle,
    FeatureStyle,
    FrameWriter,
    RowLayout,
    SlideCanvas,
    blank_slide,
    fill_background,
    run,
)

if TYPE_CHECKING:
    from pptx.presentation import Presentation

CENTER = PP_ALIGN.CENTER
MIDDLE = MSO_ANCHOR.MIDDLE
TOP = MSO_ANCHOR.TOP


def _canvas(prs: Presentation) -> SlideCanvas:
    """Add a blank slide to the deck and wrap it in a SlideCanvas."""
    return SlideCanvas(blank_slide(prs), THEME)


def slide_09_overview(prs: Presentation) -> None:
    """Workflow overview: one labeled row per phase, artifacts inline."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "4. Vue d'ensemble")
    c.title("Vue d'ensemble du workflow",
            "Du brouillon à la release taguée : phases et artefacts")
    t = TEXT
    rows = [
        ("Brouillon", BLUE, [run("Note libre → ", 10, t),
                             code_run("/process-draft"),
                             run(" → classification, renommage, branche "
                                 "d'effort", 10, t)]),
        ("Exigence", ORANGE, [code_run("/write-requirement"),
                              run(" → boucle de revue → ", 10, t),
                              code_run("/review-ask-questions"),
                              run(" + ", 10, t), code_run("/consolidate"),
                              run(" → document validé", 10, t)]),
        ("Design", BLUE, [code_run("/write-design"),
                          run(" → scope, contraintes, comportement cible, "
                              "scénarios d'acceptation → revue", 10, t)]),
        ("Plan", ORANGE, [code_run("/write-plans"),
                          run(" → étapes numérotées, tests seuils (qui "
                              "échouent d'abord), tests d'acceptation → "
                              "revue", 10, t)]),
        ("Implé.", BLUE, [code_run("/implement-step N"), run(" → ", 10, t),
                          code_run("/implementation-check N"),
                          run(" → commits groupés ", 10, t),
                          code_run("gcba")]),
        ("Release", ORANGE, [code_run("/prepare-release"),
                             run(" → merge + reword + notes → ", 10, t),
                             code_run("brel"), run(" → tag vX.Y.Z", 10, t)]),
    ]
    y0, dy = 160, 52
    for i, (label, fill, desc) in enumerate(rows):
        wf_row(c, y0 + i * dy, label, fill, desc)
    writer = c.frame(Area(50, 165 + 6 * dy, 1180, 30), anchor=MIDDLE)
    writer.add([run("pw skill", 10.5, ORANGE, bold=True),
                run(" enchaîne automatiquement les phases documentaires   "
                    "·   ", 10.5, GRAY),
                run("pw handoff", 10.5, ORANGE, bold=True),
                run(" enchaîne les étapes d'implémentation", 10.5, GRAY)],
               align=CENTER, space_after=0)
    c.footer("llm-shared - Vue d'ensemble", 9)


def slide_10_doc_phase(prs: Presentation) -> None:
    """Documentary phase: the five automated steps and the one human stop."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "5. Phase documentaire")
    c.title("Phase documentaire : automatisée",
            "pw skill enchaîne les phases sans menu ni intervention humaine "
            "(sauf revue)")
    y0, dy = 158, 66
    skill_tag = [run("== pw skill ==>", 9, ORANGE, bold=True)]
    step_row(c, y0 + 0 * dy,
             StepRow(1, BLUE, LIGHTBLUE, "/write-requirement", BLUE,
                     ": écrit le document d'exigence (feature-request ou "
                     "issue)", skill_tag))
    step_row(c, y0 + 1 * dy,
             StepRow(2, ORANGE, WARM, "/review-ask-questions", ORANGE,
                     ": l'IA pose des questions ouvertes avec options et "
                     "recommandations",
                     [run("[STOP]", 9, WHITE, bold=True)], tag_fill=ORANGE))
    step_row(c, y0 + 2 * dy,
             StepRow(3, BLUE, LIGHTBLUE,
                     "/consolidate-then-review-ask-questions", BLUE,
                     ": plie les réponses dans le document, repose des "
                     "questions si besoin",
                     [run("↻ boucle", 9, BLUE, bold=True)]))
    step_row(c, y0 + 3 * dy,
             StepRow(4, ORANGE, WARM, "/write-design", ORANGE,
                     ": scope, contraintes, comportement cible, scénarios "
                     "d'acceptation", skill_tag))
    step_row(c, y0 + 4 * dy,
             StepRow(5, BLUE, LIGHTBLUE, "/write-plans", BLUE,
                     ": plan d'exécution numéroté + plan de validation",
                     skill_tag))
    note = c.rect(Area(50, y0 + 5 * dy + 4, 1180, 44), fill=LIGHT,
                  rounded=True)
    ntf = note.text_frame
    ntf.vertical_anchor = MIDDLE
    writer = FrameWriter(ntf, THEME)
    writer.add([run("Le seul arrêt humain : répondre au tableau ", 11, BLUE),
                run("Q0x | Titre | Réponse recommandée", 11, BLUE, bold=True),
                run(" lors des revues", 11, BLUE)], align=CENTER,
               space_after=0)
    c.footer("llm-shared - Phase documentaire", 10)


def slide_11_impl_phase(prs: Presentation) -> None:
    """Implementation phase: the automated chain and its three features."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "6. Implémentation")
    c.title("Phase d'implémentation : chaîne automatisée",
            "pw handoff enchaîne les étapes sans menu ni « go-ahead » entre "
            "les cases")
    c.rect(Area(50, 158, 1180, 210), line=ORANGE, rounded=True)
    writer = c.frame(Area(60, 168, 1160, 22), anchor=MIDDLE)
    writer.add([run("ZONE AUTOMATISÉE : un seul /implement-step N déclenche "
                    "la chaîne", 11, ORANGE, bold=True)], align=CENTER,
               space_after=0)
    c.box(Area(90, 205, 500, 66), "/implement-step N",
          BoxStyle(fill=BLUE, tcolor=WHITE, size=11,
                   sub="Code + tests · ghog day vert", subcolor=TSOFT,
                   subsize=8.5))
    c.arrow(Area(595, 225, 90, 20), glyph="▼", size=16)
    c.box(Area(690, 205, 500, 66), "/implementation-check N",
          BoxStyle(fill=ORANGE, tcolor=WHITE, size=11, sub="Verdict Oui/Non",
                   subcolor=WARM, subsize=8.5))
    c.box(Area(90, 285, 500, 66), "↻ Non → /implement-missing-step",
          BoxStyle(fill=GRAY, tcolor=WHITE, size=10.5,
                   sub="Combler les écarts, ghog day, re-check",
                   subcolor=WHITE, subsize=8.5))
    c.arrow(Area(595, 305, 90, 20), glyph="ou", color=GRAY, size=12)
    c.box(Area(690, 285, 500, 66), "✓ Oui → /group-commits-msg",
          BoxStyle(fill=TEAL, tcolor=WHITE, size=10.5,
                   sub="Messages de commit groupés → a.commit",
                   subcolor=WHITE, subsize=8.5))
    feats = [
        ("🧪 Tests seuils d'abord", "Le plan commence par des tests qui "
         "échouent volontairement. Les étapes suivantes les font passer un "
         "par un.", BLUE),
        ("📦 Commits groupés", "L'IA lit le diff, groupe les fichiers du "
         "moins dépendant au plus dépendant, écrit un message par groupe.",
         ORANGE),
        ("🛑 Portail de commit", "La chaîne s'arrête à a.commit. L'auteur "
         "révise, dit « go ahead », puis gcba rejoue les commits.", BLUE),
    ]
    c.feature_row(feats, 388)
    c.footer("llm-shared - Implémentation", 11)


def slide_12_groundhog(prs: Presentation) -> None:
    """Groundhog slide: the walk chain and its three properties."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "7. Groundhog")
    c.title("Groundhog : la boucle de tests",
            "Un seul outil remplace les anciens alias : ghog day marche la "
            "chaîne et s'arrête au premier problème")
    c.box(Area(540, 158, 200, 40), "ghog day",
          BoxStyle(fill=BLUE, tcolor=WHITE, size=13))
    c.arrow(Area(630, 200, 20, 20), glyph="▼", size=16)
    steps = [
        ("1. check.bat", "Compilation + lint", LIGHTBLUE, BLUE),
        ("2. Tests affectés", "Rapide, sans couverture", WARM, ORANGE),
        ("3. Suite complète", "Couverture 100% exigée", LIGHTBLUE, BLUE),
    ]
    bw, gap = 220, 40
    total = 3 * bw + 2 * gap
    x0 = (1280 - total) / 2
    for i, (head, sub, fill, line) in enumerate(steps):
        x = x0 + i * (bw + gap)
        st = c.rect(Area(x, 228, bw, 62), fill=fill, line=line, rounded=True)
        stf = st.text_frame
        stf.vertical_anchor = MIDDLE
        writer = FrameWriter(stf, THEME)
        writer.add([run(head, 12, line, bold=True)], align=CENTER,
                   space_after=2)
        writer.add([run(sub, 9, MUTED)], align=CENTER, space_after=0)
        if i < len(steps) - 1:
            c.arrow(Area(x + bw, 245, gap, 20), glyph="→", size=16)
    c.arrow(Area(630, 292, 20, 20), glyph="▼", size=16)
    c.box(Area(415, 318, 450, 40),
          "✓ Objectif atteint : tests verts + couverture au seuil",
          BoxStyle(fill=TEAL, tcolor=WHITE, size=11))
    feats = [
        ("🔄 Boucle de correction IA", "L'IA lit le code de sortie, applique "
         "le correctif nommé dans le rapport, relance ghog day. Arrêt si "
         "aucun progrès.", BLUE),
        ("📊 Sortie maîtrisée", "Petite pour l'IA (logs redirigés, tail lu). "
         "Vivante pour l'humain (barre de progression, rapport final).",
         ORANGE),
        ("🆔 État de cycle", "Contrat fichier a.ghog.status. L'IA ne devine "
         "jamais : ghog status lit l'état, exit 6 = en cours, 7 = tué.",
         BLUE),
    ]
    c.feature_row(feats, 392, RowLayout(h=120))
    c.footer("llm-shared - Groundhog", 12)


def slide_13_commits(prs: Presentation) -> None:
    """Commits and release: the extended format and the release steps."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "8. Commits & Release")
    c.title("Commits conventionnels & préparation de release",
            "Au-delà du changelog : le « pourquoi » pour le prochain lecteur")
    banner = c.rect(Area(50, 152, 1180, 56), fill=LIGHTBLUE, rounded=True)
    btf = banner.text_frame
    btf.vertical_anchor = MIDDLE
    btf.margin_left = Pt(10)
    btf.margin_right = Pt(10)
    banner_writer = FrameWriter(btf, THEME)
    banner_writer.add([run("/group-commits-msg", 11, ORANGE, bold=True),
                       run(" : range les fichiers indexés du moins au plus "
                           "dépendant, écrit un message au format ci-dessous "
                           "par groupe dans a.commit, puis crée un commit "
                           "par groupe après relecture.", 10, BLUE)],
                      space_after=0)
    # left column: extended commit format on a dark code panel
    head_writer = c.frame(Area(50, 224, 570, 22))
    head_writer.add([run("Format de commit étendu", 12, BLUE, bold=True)],
                    space_after=0)
    panel = c.rect(Area(50, 252, 570, 240), fill=DARK, rounded=True)
    ptf = panel.text_frame
    ptf.vertical_anchor = TOP
    ptf.margin_left = Pt(12)
    ptf.margin_top = Pt(10)
    panel_writer = FrameWriter(ptf, THEME)
    panel_writer.add([run("feat", 10, ORANGE, bold=True),
                      run("(scope): sujet (52 car. max)", 10, CODE)],
                     space_after=8)
    panel_writer.add([run("Why:", 10, GREEN, bold=True)], space_after=4)
    panel_writer.add([run("Raison du changement (paragraphe 1).", 10, CODE)],
                     space_after=1)
    panel_writer.add([run("État « après » du code (paragraphe 2).", 10,
                          CODE)], space_after=8)
    panel_writer.add([run("What:", 10, GREEN, bold=True)], space_after=4)
    for m in ("- modification 1", "- modification 2", "- modification 3"):
        panel_writer.add([run(m, 10, CODE)], space_after=1)
    note_writer = c.frame(Area(50, 498, 570, 40))
    note_writer.add([run("Le titre alimente le changelog. Le corps donne le "
                         "contexte que le diff ne porte pas : pour l'humain "
                         "et pour l'IA qui relira plus tard.", 8.5, MUTED)],
                    space_after=0, line_spacing=1.1)
    # right column: numbered release steps
    rel_writer = c.frame(Area(660, 224, 570, 22))
    rel_writer.add([run("Préparation de release", 12, BLUE, bold=True)],
                   space_after=0)
    steps = [
        ("1", ORANGE, [run("/prepare-release", 10, BLUE, bold=True),
                       run(" : merge, reword, snapshot", 10, TEXT)]),
        ("2", BLUE, [run("/prepare_release_notes", 10, BLUE, bold=True),
                     run(" : titres groupés, résumé, CHANGELOG", 10, TEXT)]),
        ("3", BLUE, [run("Choix du titre de release (3 propositions)", 10,
                         TEXT)]),
        ("4", BLUE, [run("pyproject.toml + uv sync → commit chore(release)",
                         10, TEXT)]),
        ("5", ORANGE, [run("brel", 10, BLUE, bold=True),
                       run(" : build + tag vX.Y.Z [valid]", 10, TEXT)]),
    ]
    y = 256
    for num, col, desc in steps:
        o = c.oval(Area(660, y, 30, 30), col)
        otf = o.text_frame
        otf.vertical_anchor = MIDDLE
        o_writer = FrameWriter(otf, THEME)
        o_writer.add([run(num, 12, WHITE, bold=True)], align=CENTER,
                     space_after=0)
        desc_writer = c.frame(Area(700, y, 530, 30), anchor=MIDDLE)
        desc_writer.add(desc, space_after=0)
        y += 40
    tip = c.rect(Area(660, y + 4, 570, 40), fill=WARM, rounded=True)
    ttf = tip.text_frame
    ttf.vertical_anchor = MIDDLE
    tip_writer = FrameWriter(ttf, THEME)
    tip_writer.add([run("Le skill ne taggue jamais et ne pousse jamais : "
                        "c'est l'auteur qui décide", 9, ORANGE)],
                   align=CENTER, space_after=0)
    c.footer("llm-shared - Commits & Release", 13)


def slide_14_security(prs: Presentation) -> None:
    """Security slide: the two guardrails, row by row."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "Sécurité")
    c.title("🛡️ Deux garde-fous sécurité",
            "L'IA accélère le code : elle ne dispense pas des contrôles")
    # two aligned columns; left x=50 w=575, right x=655 w=575, 30px gap
    lx, rx, cw = 50, 655, 575
    left_head = c.frame(Area(lx, 158, cw, 26), anchor=MIDDLE)
    left_head.add([run("🔍 Code et dépendances", 13, BLUE, bold=True)],
                  align=CENTER, space_after=0)
    right_head = c.frame(Area(rx, 158, cw, 26), anchor=MIDDLE)
    right_head.add([run("🔒 Fuites de données", 13, ORANGE, bold=True)],
                   align=CENTER, space_after=0)
    t = TEXT
    # row 1 : Risque
    sec_cell(c, Area(lx, 192, cw, 84), "⚠️ Risque",
             [run("L'IA peut produire du code vulnérable et tirer des "
                  "bibliothèques anciennes ou trouées.", 10.5, t)],
             Sec(LIGHTBLUE, BLUE))
    sec_cell(c, Area(rx, 192, cw, 84), "⚠️ Risque",
             [run("Un secret ou une donnée métier collé dans un prompt part "
                  "chez le fournisseur du modèle.", 10.5, t)],
             Sec(WARM, ORANGE))
    # row 2 : Contrôle
    sec_cell(c, Area(lx, 288, cw, 84), "✅ Contrôle",
             [run("Passer chaque livraison au scanner ", 10.5, t),
              run("SonarQube ou équivalent", 10.5, t, bold=True),
              run(" : qualité, failles connues, dépendances.", 10.5, t)],
             Sec(LIGHTBLUE, BLUE))
    sec_cell(c, Area(rx, 288, cw, 84), "🚫 Contrôle",
             [run("Jamais de secrets ni de données réelles dans les prompts",
                  10.5, t, bold=True),
              run(" : clés, tokens, données client, extraits de prod.", 10.5,
                  t)],
             Sec(WARM, ORANGE))
    # row 3 : Porte avant prod / Détail
    sec_cell(c, Area(lx, 384, cw, 100), "🚦 Porte avant prod",
             [run("Avant la mise en production : passage SonarQube + "
                  "correction des vulnérabilités et des bibliothèques. Pas "
                  "de « go » tant que ce n'est pas propre.", 10.5, t)],
             Sec(WARM, ORANGE, accent=True))
    sec_cell(c, Area(rx, 384, cw, 100), "🧪 Détail",
             [run("On décrit le problème et la structure, on ne copie pas la "
                  "donnée ; jeux d'exemple anonymisés si besoin.", 10.5, t)],
             Sec(WARM, ORANGE))
    # row 4 : Dans le workflow
    sec_cell(c, Area(lx, 496, cw, 84), "🔗 Dans le workflow",
             [run("Après ", 10.5, t), code_run("/implementation-check"),
              run(" et la boucle ", 10.5, t), code_run("groundhog"),
              run(", avant le commit ; le vert SonarQube devient une "
                  "condition de merge.", 10.5, t)],
             Sec(LIGHTBLUE, BLUE))
    sec_cell(c, Area(rx, 496, cw, 84), "🔗 Dans le workflow",
             [run("Règle qui vaut à toutes les phases : brouillon, exigence, "
                  "design, implémentation.", 10.5, t)],
             Sec(WARM, ORANGE))
    # bottom synthesis line
    synth = c.frame(Area(50, 590, 1180, 40), anchor=MIDDLE)
    synth.add([run("💡 L'IA génère vite : ", 10.5, GRAY),
               run("SonarQube", 10.5, BLUE, bold=True),
               run(" filtre le code avant prod, la règle ", 10.5, GRAY),
               run("« jamais de secrets ni de données dans le prompt »",
                   10.5, ORANGE, bold=True),
               run(" filtre ce qui sort.", 10.5, GRAY)], align=CENTER,
              space_after=0)
    c.footer("llm-shared - Sécurité", 14)


def slide_15_benefits(prs: Presentation) -> None:
    """Benefits slide: six feature items in a two-column grid."""
    c = _canvas(prs)
    c.header(BRAND, BRAND_SUB, "Bénéfices")
    c.title("Bénéfices pour la DSI",
            "Ce que le framework apporte aux équipes et à la qualité")
    feats = [
        ("📋 Traçabilité complète", "Chaque décision est documentée. Le "
         "« pourquoi » est explicite, pas à déduire du code. Revu par l'IA "
         "avant validation.", BLUE),
        ("🤖 Multi-agent", "GitHub Copilot, Claude Code, ChatGPT Codex. Les "
         "mêmes instructions markdown, les mêmes artefacts.", ORANGE),
        ("🧪 Qualité par construction", "Tests seuils avant implémentation. "
         "Couverture 100% exigée. Revue d'implémentation systématique par "
         "l'IA.", ORANGE),
        ("📚 Documentation vivante", "Brouillon, exigence, design, plan, "
         "validation, commits : chaque phase laisse un artefact relu et "
         "consolidé.", BLUE),
        ("⚡ Automatisation", "pw skill et pw handoff enchaînent les étapes "
         "sans intervention. L'humain intervient seulement aux points de "
         "revue.", BLUE),
        ("🔓 Licence MIT", "Réutilisable dans tous les projets. Copie, fork, "
         "adaptation sans contrainte de copyleft.", ORANGE),
    ]
    w, h, gx, gy = 570, 120, 40, 20
    x0, y0 = 50, 165
    for i, (head, body, accent) in enumerate(feats):
        x = x0 + (i % 2) * (w + gx)
        y = y0 + (i // 2) * (h + gy)
        c.feature_item(Area(x, y, w, h), head, body,
                       FeatureStyle(accent=accent))
    c.footer("llm-shared - Bénéfices", 15)


def slide_16_conclusion(prs: Presentation) -> None:
    """Dark closing slide: thanks, tagline and the three stat numbers."""
    c = _canvas(prs)
    fill_background(c.slide, BLUE)
    c.rect(Area(768, 0, 512, 5), fill=ORANGE)
    brand_top_left(c)
    w1 = c.frame(Area(0, 200, 1280, 60), anchor=MIDDLE)
    w1.add([run("Merci de votre attention", 28, WHITE, bold=True)],
           align=CENTER, space_after=0)
    w2 = c.frame(Area(240, 285, 800, 70), anchor=MIDDLE)
    w2.add([run("llm-shared - du brouillon à la release", 14, TSOFT)],
           align=CENTER, space_after=4)
    w2.add([run("Sans vibe-coding, avec auto-revue IA et traçabilité "
                "complète", 11, TMUTED)], align=CENTER, space_after=0)
    stats = [("6", "phases structurées"), ("3", "agents IA compatibles"),
             ("100%", "traçabilité des décisions")]
    sw, gap = 260, 40
    total = 3 * sw + 2 * gap
    x0 = (1280 - total) / 2
    for i, (num, label) in enumerate(stats):
        x = x0 + i * (sw + gap)
        num_writer = c.frame(Area(x, 400, sw, 50), anchor=MIDDLE)
        num_writer.add([run(num, 28, ORANGE, bold=True)], align=CENTER,
                       space_after=2)
        label_writer = c.frame(Area(x, 452, sw, 24), anchor=MIDDLE)
        label_writer.add([run(label, 10, TMUTED)], align=CENTER,
                         space_after=0)


def add_main_slides(prs: Presentation) -> None:
    """Add slides 9 to 16 to the deck, in order."""
    slide_09_overview(prs)
    slide_10_doc_phase(prs)
    slide_11_impl_phase(prs)
    slide_12_groundhog(prs)
    slide_13_commits(prs)
    slide_14_security(prs)
    slide_15_benefits(prs)
    slide_16_conclusion(prs)


# eof
