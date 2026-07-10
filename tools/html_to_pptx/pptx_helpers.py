"""Primitives to rebuild an HTML slide deck as an editable PowerPoint.

The produced deck is made of text boxes and vector autoshapes only -- no
images.

Why this module exists:
    A browser-rendered HTML deck can be turned into a .pptx in two ways. One
    screenshots each slide and drops the picture on a slide: pixel-faithful
    but frozen -- nobody can edit a word in PowerPoint. The other rebuilds
    each slide from real PowerPoint objects: text frames, rounded rectangles,
    ovals, colored bars. This module is the toolbox for the second way, so
    every title, bullet and box stays editable while still matching the look
    of the source HTML.

    The trick that keeps the visual match cheap: an HTML deck authored on a
    fixed pixel canvas (here 1280x720) maps one-to-one onto a 16:9 PowerPoint
    slide (13.333in x 7.5in). 1 px == 13.333/1280 in horizontally and
    7.5/720 in vertically, and for a matched aspect ratio those two factors
    are equal. So the pixel x/y/width/height read straight out of the HTML
    become slide coordinates through Theme.inx / Theme.iny -- no re-layout,
    no guessing.

How to use it:
    Build a Theme with your brand colors and canvas size, then wrap each
    slide in a SlideCanvas and compose it from the canvas methods (frame,
    rect, oval, box, arrow) and the ready-made components (card,
    feature_item, header, title, footer). Geometry always travels as one
    Area(x, y, w, h) in source pixels, and optional looks travel in small
    style dataclasses (BoxStyle, CardStyle, ...). See
    build_llm_shared_pptx.py in this folder for a full 16-slide worked
    example, and README.md for the step-by-step method.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import TYPE_CHECKING, NamedTuple

import pptx
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from collections.abc import Sequence

    from pptx.presentation import Presentation
    from pptx.shapes.autoshape import Shape
    from pptx.slide import Slide
    from pptx.text.text import TextFrame
    from pptx.util import Length

# One styled text run: (text, size_pt, color, bold, italic). Keeping runs as
# plain tuples lets a single paragraph mix colors and weights -- for example
# a black sentence with an orange slash-command in the middle -- which is how
# the HTML inline <code> spans are reproduced.
Run = tuple[str, float, RGBColor, bool, bool]


def rgb(hex6: str) -> RGBColor:
    """Return an RGBColor from a 6-char hex string such as '00566f'."""
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def run(text: str, size: float, color: RGBColor, *, bold: bool = False,
        italic: bool = False) -> Run:
    """Describe one styled text run as a tuple for FrameWriter.add."""
    return (text, size, color, bold, italic)


class Area(NamedTuple):
    """A rectangle in source-HTML pixel coordinates: x, y, width, height."""

    x: float
    y: float
    w: float
    h: float


@dataclass
class Theme:
    """Brand colors, fonts and canvas geometry shared by every slide.

    Args:
        primary: main brand color (titles, brand name, filled boxes).
        accent: secondary color for highlights and the right accent bar.
        muted: gray used for subtitles, topic tags and secondary text.
        faint: light gray used for the footer line and page number.
        white: text color used on filled brand/accent boxes.
        font: font family applied to every run.
        canvas_px_w / canvas_px_h: pixel size of the source HTML slide.
        slide_in_w / slide_in_h: target PowerPoint slide size in inches.
        extra: free dict for any additional named colors the deck needs
            (light backgrounds, teal, dark code panel, and so on).
    """

    primary: RGBColor
    accent: RGBColor
    muted: RGBColor
    faint: RGBColor
    white: RGBColor = field(default_factory=lambda: rgb("ffffff"))
    font: str = "Segoe UI"
    canvas_px_w: int = 1280
    canvas_px_h: int = 720
    slide_in_w: float = 13.333
    slide_in_h: float = 7.5
    extra: dict[str, RGBColor] = field(default_factory=dict[str, RGBColor])

    def inx(self, px: float) -> Length:
        """Convert a horizontal pixel position/size to slide inches."""
        return Inches(px * self.slide_in_w / self.canvas_px_w)

    def iny(self, px: float) -> Length:
        """Convert a vertical pixel position/size to slide inches."""
        return Inches(px * self.slide_in_h / self.canvas_px_h)


# --- optional looks for the composed components -----------------------------
@dataclass
class BoxStyle:
    """Looks of a box(): fill, border, text color/size, optional sub line."""

    fill: RGBColor | None = None
    line: RGBColor | None = None
    tcolor: RGBColor | None = None
    size: float = 10.4
    bold: bool = True
    rounded: bool = True
    sub: str | None = None
    subcolor: RGBColor | None = None
    subsize: float = 8.0
    align: PP_ALIGN = PP_ALIGN.CENTER


@dataclass
class CardStyle:
    """Looks of a card(): accent color, optional icon, background, sizes."""

    accent: RGBColor | None = None
    icon: str | None = None
    bg: RGBColor | None = None
    title_size: float = 13
    body_size: float = 10


@dataclass
class FeatureStyle:
    """Looks of a feature_item(): accent color and text sizes."""

    accent: RGBColor | None = None
    head_size: float = 11.5
    body_size: float = 9.5


@dataclass
class RowLayout:
    """Geometry of a feature_row(): item size, gap and left origin."""

    w: float = 380
    h: float = 110
    gap: float = 20
    x0: float = 50


@dataclass
class HeaderStyle:
    """Looks of a header(): text colors and the two-tone bar split."""

    name_color: RGBColor | None = None
    sub_color: RGBColor | None = None
    topic_color: RGBColor | None = None
    split: float = 0.6


# --- presentation and slide scaffolding ------------------------------------
def make_presentation(theme: Theme) -> Presentation:
    """Return a new Presentation sized to the theme (defaults to 16:9)."""
    prs = pptx.Presentation()
    prs.slide_width = Inches(theme.slide_in_w)
    prs.slide_height = Inches(theme.slide_in_h)
    return prs


def blank_slide(prs: Presentation) -> Slide:
    """Add and return a slide using the blank layout (index 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide: Slide, color: RGBColor) -> None:
    """Paint the whole slide background one solid color (title slides)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# --- text ------------------------------------------------------------------
class FrameWriter:
    """Append styled paragraphs to one text frame.

    The writer reuses the frame's pre-existing empty paragraph for its first
    add() (a fresh text frame always starts with one paragraph), then adds a
    new paragraph per call -- the caller no longer tracks a `first` flag.
    """

    def __init__(self, tf: TextFrame, theme: Theme) -> None:
        """Wrap the text frame and remember the theme for the font family."""
        self._tf = tf
        self._theme = theme
        self._used = False

    def add(self, runs: Sequence[Run], *, align: PP_ALIGN = PP_ALIGN.LEFT,
            space_after: float = 2, line_spacing: float = 1.0) -> None:
        """Append one paragraph made of `runs` to the frame."""
        p = self._tf.add_paragraph() if self._used else self._tf.paragraphs[0]
        self._used = True
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for text, size, color, bold, italic in runs:
            r = p.add_run()
            r.text = text
            f = r.font
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color
            f.name = self._theme.font


# --- per-slide drawing surface ----------------------------------------------
class SlideCanvas:
    """One slide plus the theme: every drawing primitive hangs off this.

    Methods take their geometry as one Area in source-HTML pixels and their
    optional looks as a small style dataclass, so each call stays short and
    the pixel numbers read straight out of the HTML source.
    """

    def __init__(self, slide: Slide, theme: Theme) -> None:
        """Bind the drawing surface to one slide and one theme."""
        self.slide = slide
        self.theme = theme

    # --- text ---------------------------------------------------------------
    def frame(self, area: Area, *, anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> FrameWriter:
        """Add an empty textbox at pixel coordinates and return its writer."""
        t = self.theme
        tb = self.slide.shapes.add_textbox(t.inx(area.x), t.iny(area.y),
                                           t.inx(area.w), t.iny(area.h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        return FrameWriter(tf, t)

    def shape_frame(self, shape: Shape, *, anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
                    pad: float = 4) -> FrameWriter:
        """Prepare a shape's own text frame and return its writer.

        Centering vertically (MIDDLE) is the default because most boxes hold
        one or two short centered lines.
        """
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Pt(pad)
        tf.margin_right = Pt(pad)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        return FrameWriter(tf, self.theme)

    # --- shapes ---------------------------------------------------------------
    def rect(self, area: Area, *, fill: RGBColor | None = None,
             line: RGBColor | None = None, rounded: bool = False,
             line_w: float = 1.5) -> Shape:
        """Add a rectangle (optionally rounded) at pixel coordinates.

        Pass fill=None for a transparent body and line=None for no border, so
        the same helper draws a filled chip, an outlined box, or a thin
        colored bar.
        """
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        t = self.theme
        s = self.slide.shapes.add_shape(kind, t.inx(area.x), t.iny(area.y),
                                        t.inx(area.w), t.iny(area.h))
        if fill is None:
            s.fill.background()
        else:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        if line is None:
            s.line.fill.background()
        else:
            s.line.color.rgb = line
            s.line.width = Pt(line_w)
        s.shadow.inherit = False
        return s

    def oval(self, area: Area, fill: RGBColor, *,
             line: RGBColor | None = None) -> Shape:
        """Add a filled circle/ellipse, used for numbered step badges."""
        t = self.theme
        s = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, t.inx(area.x),
                                        t.iny(area.y), t.inx(area.w),
                                        t.iny(area.h))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if line is None:
            s.line.fill.background()
        else:
            s.line.color.rgb = line
        s.shadow.inherit = False
        return s

    def box(self, area: Area, main: str, style: BoxStyle | None = None) -> Shape:
        """Draw a rounded box with a centered main line and optional sub line.

        This is the workhorse for flow diagrams: a filled brand box, an
        outlined box, or a two-line box (a label plus a smaller caption
        underneath).
        """
        style = style if style is not None else BoxStyle()
        tcolor = style.tcolor if style.tcolor is not None else self.theme.primary
        s = self.rect(area, fill=style.fill, line=style.line,
                      rounded=style.rounded)
        writer = self.shape_frame(s)
        writer.add([run(main, style.size, tcolor, bold=style.bold)],
                   align=style.align, space_after=0)
        if style.sub is not None:
            subcolor = style.subcolor if style.subcolor is not None else tcolor
            writer.add([run(style.sub, style.subsize, subcolor)],
                       align=style.align, space_after=0)
        return s

    def arrow(self, area: Area, *, glyph: str = "▼",
              color: RGBColor | None = None, size: float = 13) -> None:
        """Place a small centered connector glyph (▼, →, 'ou', ...)."""
        color = color if color is not None else self.theme.muted
        writer = self.frame(area, anchor=MSO_ANCHOR.MIDDLE)
        writer.add([run(glyph, size, color)], align=PP_ALIGN.CENTER,
                   space_after=0)

    # --- reusable components ---------------------------------------------------
    def card(self, area: Area, ctitle: str, body: str,
             style: CardStyle | None = None) -> None:
        """A rounded panel with a top accent bar, a colored title and body.

        Mirrors the HTML .card component (top border in the accent color).
        """
        style = style if style is not None else CardStyle()
        theme = self.theme
        accent = style.accent if style.accent is not None else theme.primary
        bg = style.bg if style.bg is not None else theme.extra.get(
            "light", rgb("f0f4f6"))
        self.rect(area, fill=bg, rounded=True)
        self.rect(Area(area.x, area.y, area.w, 6), fill=accent)
        writer = self.frame(Area(area.x + 16, area.y + 16, area.w - 32,
                                 area.h - 28))
        if style.icon:
            writer.add([run(style.icon, 24,
                            theme.extra.get("text", rgb("333333")))],
                       space_after=4)
        writer.add([run(ctitle, style.title_size, accent, bold=True)],
                   space_after=6)
        writer.add([run(body, style.body_size,
                        theme.extra.get("muted2", rgb("555555")))],
                   space_after=0, line_spacing=1.15)

    def feature_item(self, area: Area, head: str, body: str,
                     style: FeatureStyle | None = None) -> None:
        """A rounded panel with a left accent bar, a heading and body text.

        Mirrors the HTML .feature-item component (left border in the accent).
        """
        style = style if style is not None else FeatureStyle()
        theme = self.theme
        accent = style.accent if style.accent is not None else theme.primary
        bg = theme.extra.get("light", rgb("f0f4f6"))
        self.rect(area, fill=bg, rounded=True)
        self.rect(Area(area.x, area.y, 6, area.h), fill=accent)
        writer = self.frame(Area(area.x + 16, area.y + 12, area.w - 26,
                                 area.h - 20))
        writer.add([run(head, style.head_size, theme.primary, bold=True)],
                   space_after=4)
        writer.add([run(body, style.body_size,
                        theme.extra.get("muted", rgb("666666")))],
                   space_after=0, line_spacing=1.15)

    def feature_row(self, feats: Sequence[tuple[str, str, RGBColor]], y: float,
                    layout: RowLayout | None = None) -> None:
        """Lay a list of (head, body, accent) feature items in one row."""
        layout = layout if layout is not None else RowLayout()
        for i, (head, body, accent) in enumerate(feats):
            self.feature_item(
                Area(layout.x0 + i * (layout.w + layout.gap), y, layout.w,
                     layout.h),
                head, body, FeatureStyle(accent=accent))

    # --- shared chrome (header band, title block, footer) ----------------------
    def header(self, brand_name: str, brand_sub: str, topic: str,
               style: HeaderStyle | None = None) -> None:
        """Draw the top brand line, the topic tag and the two-tone accent bar.

        The accent bar is two rectangles: primary on the left `split`
        fraction, accent on the remainder -- the same 60/40 blue/orange
        split as the HTML.
        """
        style = style if style is not None else HeaderStyle()
        theme = self.theme
        name_color = style.name_color or theme.primary
        sub_color = style.sub_color or theme.muted
        topic_color = style.topic_color or theme.muted
        writer = self.frame(Area(40, 12, 500, 40))
        writer.add([run(brand_name, 13, name_color, bold=True)],
                   space_after=1)
        writer.add([run(brand_sub, 7, sub_color)], space_after=0)
        writer2 = self.frame(Area(780, 20, 460, 20))
        writer2.add([run(topic, 8.5, topic_color)], align=PP_ALIGN.RIGHT,
                    space_after=0)
        split_px = theme.canvas_px_w * style.split
        self.rect(Area(0, 58, split_px, 5), fill=theme.primary)
        self.rect(Area(split_px, 58, theme.canvas_px_w - split_px, 5),
                  fill=theme.accent)

    def title(self, t: str, sub: str) -> None:
        """Draw the slide title and subtitle block below the accent bar."""
        theme = self.theme
        writer = self.frame(Area(50, 72, 1180, 44))
        writer.add([run(t, 21, theme.primary, bold=True)], space_after=0)
        writer2 = self.frame(Area(50, 116, 1180, 28))
        writer2.add([run(sub, 10.5, theme.muted)], space_after=0)

    def footer(self, left: str, num: int) -> None:
        """Draw the bottom-left caption and the bottom-right page number."""
        theme = self.theme
        writer = self.frame(Area(40, 700, 500, 16))
        writer.add([run(left, 7.5, theme.faint)], space_after=0)
        writer2 = self.frame(Area(1150, 700, 90, 16))
        writer2.add([run(str(num), 7.5, theme.faint)], align=PP_ALIGN.RIGHT,
                    space_after=0)


# eof
