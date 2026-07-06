"""Worked example: rebuild docs/llm-shared_presentation.html as a native,
editable PowerPoint using the primitives in pptx_helpers.py.

Run it from anywhere:

    python tools/html_to_pptx/build_llm_shared_pptx.py

It writes docs/llm-shared_presentation.pptx (16 slides, text and vector
shapes only -- no images). Every pixel coordinate below is read from the
matching slide in the HTML source; the Theme maps 1280x720 px onto a
13.333in x 7.5in 16:9 slide, so the numbers transfer without re-layout.

Treat this file as the pattern to copy for any other HTML deck: keep the
generic drawing in pptx_helpers.py, and keep one build_<deck>_pptx.py like
this one per presentation, holding the palette, the per-slide coordinates
and the deck-specific compositions (the flow columns and step rows below).
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx.util import Pt                                          # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                   # noqa: E402
from pptx.enum.dml import MSO_LINE                                # noqa: E402

from pptx_helpers import (                                        # noqa: E402
    Theme, rgb, make_presentation, blank_slide, fill_background,
    run, add_paragraph, textbox, rect, oval, box, arrow, card,
    feature_item, feature_row, header, title, footer,
)

CENTER, LEFT, RIGHT = PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT
MIDDLE, TOP = MSO_ANCHOR.MIDDLE, MSO_ANCHOR.TOP

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.abspath(os.path.join(HERE, "..", "..", "docs",
                                   "llm-shared_presentation.pptx"))

# --- palette (the HTML :root variables) ------------------------------------
BLUE = rgb("00566f")
ORANGE = rgb("ec6608")
GREEN = rgb("95c11f")
GRAY = rgb("717e86")
LGRAY = rgb("aaaaaa")
LIGHTBLUE = rgb("e6f0f3")
WARM = rgb("fff3e6")
LIGHT = rgb("f0f4f6")
TEAL = rgb("006d67")
DARK = rgb("1a2a30")
WHITE = rgb("ffffff")
TEXT = rgb("333333")
MUTED = rgb("666666")
MUTED2 = rgb("555555")
TSOFT = rgb("c4d8de")
TMUTED = rgb("8eb2bb")
CODE = rgb("8eb2bb")

THEME = Theme(primary=BLUE, accent=ORANGE, muted=GRAY, faint=LGRAY,
              white=WHITE, extra={"light": LIGHT, "muted": MUTED,
                                  "muted2": MUTED2, "text": TEXT})

# Brand placeholders. Like the HTML deck (which swaps them at render time
# through llm-shared_presentation.local.js), the build reads LLM_SHARED_BRAND
# and LLM_SHARED_BRAND_SUB from the environment and keeps the placeholders
# when they are not set.
BRAND = os.environ.get("LLM_SHARED_BRAND", "Organization name")
BRAND_SUB = os.environ.get("LLM_SHARED_BRAND_SUB", "ORGANIZATION SUBTITLE")


# --- deck-specific compositions (copy and adapt these per deck) -------------
def phase_box(s, x, y, w, h, num, label, detail, style):
    """One phase card from slide 4: number, label, detail, three fill styles."""
    if style == "active":
        fill, line, tc, dc = BLUE, BLUE, WHITE, TSOFT
    elif style == "active-orange":
        fill, line, tc, dc = ORANGE, ORANGE, WHITE, WARM
    else:
        fill, line, tc, dc = WHITE, BLUE, BLUE, MUTED
    rect(s, THEME, x, y, w, h, fill=fill, line=line, rounded=True)
    tf = textbox(s, THEME, x + 4, y + 12, w - 8, h - 20)
    add_paragraph(tf, [run(str(num), 24, ORANGE if style == "plain" else
                           WHITE, True)], THEME, align=CENTER, first=True,
                  space_after=2)
    add_paragraph(tf, [run(label, 12, tc, True)], THEME, align=CENTER,
                  space_after=3)
    add_paragraph(tf, [run(detail, 8.5, dc)], THEME, align=CENTER,
                  space_after=0, line_spacing=1.05)


def loop_col(s, cx, y0, title_txt, sub_txt, steps, box_w=240, box_h=42,
             gap=8):
    """A vertical flow column of loop boxes with ▼ connectors and notes."""
    tf = textbox(s, THEME, cx - 220, y0 - 60, 440, 24, anchor=MIDDLE)
    add_paragraph(tf, [run(title_txt, 13, BLUE, True)], THEME, align=CENTER,
                  first=True, space_after=0)
    if sub_txt:
        tf = textbox(s, THEME, cx - 240, y0 - 36, 480, 22, anchor=MIDDLE)
        add_paragraph(tf, sub_txt, THEME, align=CENTER, first=True,
                      space_after=0)
    y = y0
    for i, st in enumerate(steps):
        style = st["style"]
        if style == "filled":
            fill, line, tc = BLUE, BLUE, WHITE
        elif style == "orange":
            fill, line, tc = WHITE, ORANGE, ORANGE
        elif style == "filled-orange":
            fill, line, tc = ORANGE, ORANGE, WHITE
        else:
            fill, line, tc = WHITE, BLUE, BLUE
        box(s, THEME, cx - box_w / 2, y, box_w, box_h, st["text"], fill=fill,
            line=line, tcolor=tc, size=10, bold=True)
        y += box_h
        if st.get("note"):
            tf = textbox(s, THEME, cx - box_w / 2, y + 1, box_w, 16,
                         anchor=MIDDLE)
            add_paragraph(tf, [run(st["note"], 8.5, st.get("note_color",
                          ORANGE), True)], THEME, align=CENTER, first=True,
                          space_after=0)
            y += 17
        if i < len(steps) - 1:
            arrow(s, THEME, cx - 10, y, 20)
            y += 20


def code_run(text):
    """An orange slash-command run, matching the HTML inline <code> spans."""
    return run(text, 10, ORANGE, False)


def wf_row(s, y, label, label_fill, desc_runs):
    """One workflow row from slide 7: colored label + arrow + tinted body."""
    lab = rect(s, THEME, 50, y, 120, 40, fill=label_fill, rounded=True)
    ltf = lab.text_frame
    ltf.vertical_anchor = MIDDLE
    ltf.word_wrap = True
    add_paragraph(ltf, [run(label, 11, WHITE, True)], THEME, align=CENTER,
                  first=True, space_after=0)
    arrow(s, THEME, 172, y + 10, 24, glyph="→", size=14)
    bg = LIGHTBLUE if label_fill == BLUE else WARM
    body = rect(s, THEME, 200, y, 1030, 40, fill=bg, rounded=True)
    btf = body.text_frame
    btf.vertical_anchor = MIDDLE
    btf.word_wrap = True
    btf.margin_left = Pt(8)
    add_paragraph(btf, desc_runs, THEME, first=True, space_after=0)


def step_row(s, y, num, circle_color, box_fill, cmd, cmd_color, desc,
             tag_runs, tag_fill=None):
    """One numbered step row from slide 8: oval badge + command box + tag."""
    o = oval(s, THEME, 50, y, 44, 44, fill=circle_color)
    otf = o.text_frame
    otf.vertical_anchor = MIDDLE
    add_paragraph(otf, [run(str(num), 16, WHITE, True)], THEME, align=CENTER,
                  first=True, space_after=0)
    b = rect(s, THEME, 110, y, 940, 44, fill=box_fill, rounded=True)
    btf = b.text_frame
    btf.vertical_anchor = MIDDLE
    btf.word_wrap = True
    btf.margin_left = Pt(8)
    add_paragraph(btf, [run(cmd, 12, cmd_color, True),
                        run("  " + desc, 10, MUTED)], THEME, first=True,
                  space_after=0)
    if tag_fill is not None:
        t = rect(s, THEME, 1075, y + 6, 155, 32, fill=tag_fill, rounded=True)
        ttf = t.text_frame
        ttf.vertical_anchor = MIDDLE
        add_paragraph(ttf, tag_runs, THEME, align=CENTER, first=True,
                      space_after=0)
    else:
        tf = textbox(s, THEME, 1075, y, 155, 44, anchor=MIDDLE)
        add_paragraph(tf, tag_runs, THEME, align=CENTER, first=True,
                      space_after=0)


def brand_top_left(s):
    """White brand line for the two dark title slides (1 and 16)."""
    tf = textbox(s, THEME, 40, 30, 500, 44)
    add_paragraph(tf, [run(BRAND, 13, WHITE, True)], THEME, first=True,
                  space_after=1)
    add_paragraph(tf, [run(BRAND_SUB, 7, TMUTED)], THEME, space_after=0)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_01_title(prs):
    s = blank_slide(prs)
    fill_background(s, BLUE)
    rect(s, THEME, 768, 0, 512, 5, fill=ORANGE)
    brand_top_left(s)
    tf = textbox(s, THEME, 0, 250, 1280, 70, anchor=MIDDLE)
    add_paragraph(tf, [run("llm-shared", 42, WHITE, True)], THEME,
                  align=CENTER, first=True, space_after=0)
    tf = textbox(s, THEME, 240, 340, 800, 70, anchor=MIDDLE)
    add_paragraph(tf, [run("Framework de développement assisté par IA", 16,
                  TSOFT)], THEME, align=CENTER, first=True, space_after=2)
    add_paragraph(tf, [run("De l'idée brute à la release taguée", 16, TSOFT)],
                  THEME, align=CENTER, space_after=0)
    tf = textbox(s, THEME, 0, 450, 1280, 30, anchor=MIDDLE)
    add_paragraph(tf, [run("Direction Informatique - 2025", 13, TMUTED)],
                  THEME, align=CENTER, first=True, space_after=0)


def slide_02_agenda(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "llm-shared")
    title(s, THEME, "Agenda",
          "Framework de développement structuré avec l'IA")
    items = [
        "Le problème du « vibe-coding »",
        "La solution : un workflow structuré",
        "Principe clé : l'auto-revue IA",
        "Vue d'ensemble du workflow",
        "Phase documentaire",
        "Phase d'implémentation",
        "Commits conventionnels & release",
        "Groundhog & automatisation",
        "Sécurité : deux garde-fous",
    ]
    col_x = [150, 700]
    y0, dy = 190, 78
    for i, text in enumerate(items):
        x = col_x[i % 2]
        y = y0 + (i // 2) * dy
        tf = textbox(s, THEME, x, y, 50, 44, anchor=MIDDLE)
        add_paragraph(tf, [run(str(i + 1), 26, ORANGE, True)], THEME,
                      first=True, space_after=0)
        tf = textbox(s, THEME, x + 55, y, 420, 44, anchor=MIDDLE)
        add_paragraph(tf, [run(text, 13, BLUE, True)], THEME, first=True,
                      space_after=0)
    footer(s, THEME, "llm-shared - Agenda", 2)


def slide_03_problem(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "1. Problème")
    title(s, THEME, "Le problème du « vibe-coding »",
          "Générer du code vite sans spécifier ce qu'on veut vraiment")
    cards = [
        ("💡", "L'idée",
         "« J'ai une idée, génère-moi du code, je verrai les détails plus "
         "tard. »"),
        ("⚡", "Le code arrive vite",
         "Mais il s'appuie sur des hypothèses non écrites. Il casse au "
         "premier cas réel non prévu."),
        ("💥", "La conséquence",
         "Code sans documentation d'intention. Difficile à maintenir, à "
         "expliquer, à étendre : même par l'auteur."),
    ]
    x0, w, gap = 50, 380, 20
    for i, (icon, ct, body) in enumerate(cards):
        card(s, THEME, x0 + i * (w + gap), 165, w, 190, ct, body,
             accent=ORANGE, icon=icon)
    rect(s, THEME, 50, 385, 1180, 115, fill=WARM, line=ORANGE, rounded=True)
    tf = textbox(s, THEME, 60, 400, 1160, 95, anchor=MIDDLE)
    add_paragraph(tf, [run("⚠️ Ce que refuse le workflow", 14, ORANGE, True)],
                  THEME, align=CENTER, first=True, space_after=6)
    add_paragraph(tf, [run("Passer directement à l'implémentation sans "
                  "exigence écrite · Sans design · Sans tests · Sans "
                  "scénarios d'acceptation", 11, MUTED2)], THEME,
                  align=CENTER, space_after=0)
    footer(s, THEME, "llm-shared - Problème", 3)


def slide_04_solution(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "2. Solution")
    title(s, THEME, "La solution : un workflow en phases",
          "Garder la créativité, refuser le raccourci vers le code")
    phases = [
        ("1", "Brouillon", "Note libre en langage naturel", "plain"),
        ("2", "Exigence", "Document structuré + revue", "active-orange"),
        ("3", "Design", "Solution + scénarios d'acceptation", "active"),
        ("4", "Plan", "Étapes numérotées + tests", "active-orange"),
        ("5", "Implémentation", "Code + tests + auto-revue", "active"),
        ("6", "Release", "Notes + tag de version", "active-orange"),
    ]
    w, gap = 178, 18
    total = len(phases) * w + (len(phases) - 1) * gap
    x0 = (1280 - total) / 2
    y = 165
    for i, (num, label, detail, style) in enumerate(phases):
        x = x0 + i * (w + gap)
        phase_box(s, x, y, w, 150, num, label, detail, style)
        if i < len(phases) - 1:
            arrow(s, THEME, x + w, y + 55, gap, glyph="→", size=16)
    card(s, THEME, 50, 355, 580, 175, "Chaque phase laisse une trace",
         "Brouillon, exigence, design, plan, validation, commits groupés : "
         "le workflow double comme documentation du processus de "
         "développement.", accent=BLUE)
    card(s, THEME, 650, 355, 580, 175, "Compatible avec tout LLM",
         "Les commandes slash se résolvent en markdown brut. Copilot, Claude "
         "Code, ChatGPT Codex : tout agent qui lit et écrit des fichiers peut "
         "suivre le workflow.", accent=ORANGE)
    footer(s, THEME, "llm-shared - Solution", 4)


def slide_05_review(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "3. Auto-revue")
    title(s, THEME, "Principe clé : l'IA révise son propre travail",
          "Ne jamais faire confiance au premier passage : même pour la "
          "documentation")
    left = [
        {"text": "Écriture du document", "style": "filled"},
        {"text": "IA pose des questions ouvertes", "style": "orange",
         "note": "[STOP] Humain répond", "note_color": ORANGE},
        {"text": "Consolidation des réponses", "style": "plain",
         "note": "↻ Nouvelles questions ? Boucle", "note_color": BLUE},
        {"text": "Document validé ✓", "style": "filled-orange"},
    ]
    right = [
        {"text": "Génération du code", "style": "filled"},
        {"text": "Vérification : code = plan ?", "style": "orange",
         "note": "[STOP] Écarts détectés ?", "note_color": ORANGE},
        {"text": "Corrections itératives", "style": "plain",
         "note": "↻ Vérification jusqu'à conformité", "note_color": BLUE},
        {"text": "Étape validée ✓", "style": "filled-orange"},
    ]
    loop_col(s, 340, 230, "Revue documentaire", None, left)
    loop_col(s, 940, 230, "Revue d'implémentation", None, right)
    footer(s, THEME, "llm-shared - Auto-revue", 5)


def slide_06_review_practice(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "3. Auto-revue")
    title(s, THEME, "L'auto-revue en pratique : deux skills dédiés",
          "/review-ask-questions relit le document, /implementation-check "
          "relit le code")
    left = [
        {"text": "Question : contexte & problème", "style": "filled"},
        {"text": "Reformulation BBQ des enjeux", "style": "plain"},
        {"text": "Options X1 / X2 / X3 : pour & contre", "style": "orange"},
        {"text": "Option recommandée + arguments", "style": "plain",
         "note": "[STOP] Humain répond", "note_color": ORANGE},
        {"text": "Réponse : option XY + tableau récap ✓",
         "style": "filled-orange"},
    ]
    right = [
        {"text": "Verdict : Oui / Non - étape N", "style": "filled"},
        {"text": "Ce qui est fait / travail manquant", "style": "orange",
         "note": "[STOP] Écarts → implement-missing", "note_color": ORANGE},
        {"text": "Architecture DDD-Hexagonal", "style": "plain"},
        {"text": "Performance : pas de O(n²)", "style": "plain"},
        {"text": "Couverture unitaire 100% ?", "style": "plain"},
        {"text": "Intégrité fonctionnelle préservée ✓",
         "style": "filled-orange"},
    ]
    sub_l = [run("l'IA relit ce qu'elle vient de rédiger : ", 9, GRAY, False,
                 True), run("/review-ask-questions", 9, ORANGE, False, True)]
    sub_r = [run("l'IA relit ce qu'elle vient de coder : ", 9, GRAY, False,
                 True), run("/implementation-check", 9, ORANGE, False, True)]
    loop_col(s, 340, 235, "Revue documentaire", sub_l, left, box_h=38, gap=6)
    loop_col(s, 940, 235, "Revue d'implémentation", sub_r, right, box_h=38,
             gap=6)
    footer(s, THEME, "llm-shared - Auto-revue en pratique", 6)


def ex_cell(s, x, y, w, h, fill, hcolor, heading, paras, line=None,
            accent=None, dashed=False, hsize=11.5, bsize=10):
    """One example panel from slides 7-8: tinted rounded box, colored bold
    heading, then one body paragraph per entry in paras (lists of runs)."""
    r = rect(s, THEME, x, y, w, h, fill=fill, line=line, rounded=True)
    if dashed:
        r.line.dash_style = MSO_LINE.DASH
    if accent is not None:
        rect(s, THEME, x, y, 5, h, fill=accent)
    tf = textbox(s, THEME, x + 12, y + 8, w - 22, h - 14)
    add_paragraph(tf, [run(heading, hsize, hcolor, True)], THEME, first=True,
                  space_after=3)
    for pr in paras:
        add_paragraph(tf, pr, THEME, space_after=3, line_spacing=1.1)


def ex_banner(s, y, h, badge, fill, line, tcolor, sentence, note_runs=None):
    """Full-width AVANT/APRÈS sentence band with a bold badge run."""
    b = rect(s, THEME, 50, y, 1180, h, fill=fill, line=line, rounded=True)
    tf = b.text_frame
    tf.vertical_anchor = MIDDLE
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    runs = [run(badge + "   ", 11, tcolor, True),
            run(sentence, 11, tcolor, False, True)]
    if note_runs:
        runs += note_runs
    add_paragraph(tf, runs, THEME, first=True, space_after=0,
                  line_spacing=1.1)


def ex_connector(s, y, runs):
    """Centered ▼ connector line between the example stages."""
    tf = textbox(s, THEME, 50, y, 1180, 20, anchor=MIDDLE)
    add_paragraph(tf, [run("▼  ", 10, GRAY)] + runs, THEME, align=CENTER,
                  first=True, space_after=0)


def slide_07_doc_example(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "3. Auto-revue")
    title(s, THEME, "Auto-revue documentaire : un exemple réel",
          "my-project v9.1.0 « date validation » : de la phrase ambiguë à "
          "la décision tracée")
    ex_banner(s, 152, 52, "AVANT", WARM, ORANGE, ORANGE,
              '"If the date in the filename is invalid, the file is renamed '
              '_DATE_INCOHERENTE and not processed."',
              [run("  — « invalid » cache deux cas différents.", 11,
                   ORANGE)])
    ex_connector(s, 208, [run("/review-ask-questions", 10, ORANGE, True),
                 run(" : l'IA détecte l'ambiguïté et pose Q01", 10, GRAY)])
    # Q01: the wordy phrasing next to the readable one
    ex_cell(s, 50, 232, 580, 80, LIGHTBLUE, BLUE, "Q01 — version verbeuse",
            [[run('"Should an eight-digit token failing calendar-level '
                  'plausibility verification be assimilated to the '
                  'incoherence semantics of the ten-year window rule…?"',
                  10, MUTED2)]])
    ex_cell(s, 650, 232, 580, 80, LIGHTBLUE, BLUE, "Q01 — la même, lisible",
            [[run('"When the token is not a real date (like 32132024): '
                  'block the file, or just use today\'s date?"', 10,
                  MUTED2)]], accent=ORANGE)
    # options with pros and cons; A2 carries the recommendation border
    ex_cell(s, 50, 318, 380, 84, LIGHT, BLUE, "A1 — Block the file",
            [[run("✓ operator sees every odd token", 10, MUTED2)],
             [run("✗ goes beyond the CDC rule", 10, MUTED2)]])
    ex_cell(s, 450, 318, 380, 84, LIGHTBLUE, BLUE,
            "A2 — Use today's date ★",
            [[run("✓ matches the CDC, backward compatible", 10, MUTED2)],
             [run("✗ a typo is processed silently", 10, MUTED2)]],
            line=BLUE)
    ex_cell(s, 850, 318, 380, 84, LIGHT, BLUE, "A3 — Fuzzy typo filter",
            [[run("✓ catches likely typos", 10, MUTED2)],
             [run("✗ arbitrary, hard to test", 10, MUTED2)]])
    # the pre-filled answer and the free human choice
    ex_cell(s, 50, 408, 580, 72, LIGHT, BLUE, "Réponse pré-remplie par l'IA",
            [[run('"Answer to Q01: option A2 — matches the CDC scope and '
                  'the behavior in production."', 10, MUTED2)]])
    ex_cell(s, 650, 408, 580, 72, WARM, ORANGE,
            "[STOP] Réponse humaine — libre",
            [[run('"Answer to Q01: A4 — today\'s date, but log a WARNING." ',
                  10, MUTED2),
              run("Un 4e choix est permis.", 10, MUTED2)]])
    ex_connector(s, 484, [run("/consolidate-then-review-ask-questions", 10,
                 ORANGE, True), run(" : la réponse est pliée dans le "
                 "document", 10, GRAY)])
    # one-row decision table
    rect(s, THEME, 50, 508, 1180, 62, fill=LIGHT, rounded=True)
    cols = [
        (66, 50, "ID", "Q01"),
        (126, 360, "Question", "Impossible date: incoherent or missing?"),
        (496, 250, "Chosen option", "A2 + WARNING (A4)"),
        (756, 460, "Rationale",
         'CDC scopes "incoherent" to the 10-year window.'),
    ]
    for x, w, head, value in cols:
        tf = textbox(s, THEME, x, 514, w, 50)
        add_paragraph(tf, [run(head, 10, BLUE, True)], THEME, first=True,
                      space_after=2)
        add_paragraph(tf, [run(value, 9.5, MUTED2)], THEME, space_after=0)
    ex_banner(s, 576, 56, "APRÈS", TEAL, None, WHITE,
              '"A real date outside [today - 10 years, today] → file '
              'renamed, not processed. A token that is not a real date → '
              'missing date: the process starts with today\'s date."')
    footer(s, THEME, "llm-shared - Auto-revue : exemple documentaire", 7)


def slide_08_impl_example(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "3. Auto-revue")
    title(s, THEME, "Auto-revue d'implémentation : un exemple réel",
          "my-project v9.11.0 : le plan disait « fait », la revue du "
          "comportement réel a trouvé le morceau manquant")
    ex_cell(s, 50, 170, 580, 110, LIGHTBLUE, BLUE,
            "🎯 Ce que le plan demandait",
            [[run("Show the preparation progress in two places: the live "
                  "widget on screen, and the progress log file that "
                  "operators read.", 11, MUTED2)]], hsize=13)
    ex_cell(s, 650, 170, 580, 110, WARM, ORANGE,
            "AVANT — Implémenté au premier jet",
            [[run("The widget shows the progress, live. The log file is "
                  "ready to display it too. All tests are green. ✓", 11,
                  MUTED2)]], line=ORANGE, hsize=13)
    ex_connector(s, 292, [run("/implementation-check", 10, ORANGE, True),
                 run(" + revue du comportement réel → ", 10, GRAY),
                 run("verdict : Non", 10, ORANGE, True)])
    ex_cell(s, 50, 320, 580, 200, WHITE, ORANGE,
            "🔍 L'écart détecté par la revue",
            [[run("The log file never receives the progress. The widget "
                  "moves forward, but the file stays frozen on an old "
                  "value. The display side was built — the part that "
                  "writes the progress into the file was not.", 10.5,
                  MUTED2)],
             [run("Why did it slip through? ", 10.5, ORANGE, True),
              run("The AI built each half separately — and tested each "
                  "half separately: the widget with live data, the log "
                  "display with a hand-written file. Both halves pass. The "
                  "wire between them was never written, so no test could "
                  "fail.", 10.5, MUTED2)]], line=ORANGE, dashed=True,
            hsize=13)
    ex_cell(s, 650, 320, 580, 200, LIGHTBLUE, BLUE, "🔧 Le correctif",
            [[run("One fix commit: recognize the progress signal during "
                  "preparation, pass the counters to the file before each "
                  "write, and keep a preparation update from blocking the "
                  "next splitting update.", 10.5, MUTED2)],
             [run("Plus the missing test: ", 10.5, BLUE, True),
              run("a new test now follows the progress end to end, from "
                  "the running phase to the file content — the gap can no "
                  "longer come back.", 10.5, MUTED2)]], hsize=13)
    ex_connector(s, 532, [run("nouvelle passe ", 10, GRAY),
                 run("/implementation-check", 10, ORANGE, True),
                 run(" → ", 10, GRAY), run("verdict : Oui", 10, TEAL, True)])
    b = rect(s, THEME, 50, 560, 1180, 70, fill=TEAL, rounded=True)
    tf = b.text_frame
    tf.vertical_anchor = MIDDLE
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    add_paragraph(tf, [run("APRÈS   ", 11, WHITE, True),
                  run("The log file now moves at the same pace as the "
                      "widget: operators see the same progress in both "
                      "places. Three new tests lock the behavior — full "
                      "suite green, coverage 100%.", 11, WHITE)], THEME,
                  first=True, space_after=0, line_spacing=1.1)
    footer(s, THEME, "llm-shared - Auto-revue : exemple implémentation", 8)


def slide_09_overview(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "4. Vue d'ensemble")
    title(s, THEME, "Vue d'ensemble du workflow",
          "Du brouillon à la release taguée : phases et artefacts")
    t = TEXT
    rows = [
        ("Brouillon", BLUE, [run("Note libre → ", 10, t), code_run(
            "/process-draft"), run(" → classification, renommage, branche "
                                   "d'effort", 10, t)]),
        ("Exigence", ORANGE, [code_run("/write-requirement"), run(
            " → boucle de revue → ", 10, t), code_run("/review-ask-questions"),
            run(" + ", 10, t), code_run("/consolidate"),
            run(" → document validé", 10, t)]),
        ("Design", BLUE, [code_run("/write-design"), run(
            " → scope, contraintes, comportement cible, scénarios "
            "d'acceptation → revue", 10, t)]),
        ("Plan", ORANGE, [code_run("/write-plans"), run(
            " → étapes numérotées, tests seuils (qui échouent d'abord), tests "
            "d'acceptation → revue", 10, t)]),
        ("Implé.", BLUE, [code_run("/implement-step N"), run(" → ", 10, t),
            code_run("/implementation-check N"), run(" → commits groupés ", 10,
            t), code_run("gcba")]),
        ("Release", ORANGE, [code_run("/prepare-release"), run(
            " → merge + reword + notes → ", 10, t), code_run("brel"),
            run(" → tag vX.Y.Z", 10, t)]),
    ]
    y0, dy = 160, 52
    for i, (label, fill, desc) in enumerate(rows):
        wf_row(s, y0 + i * dy, label, fill, desc)
    tf = textbox(s, THEME, 50, 165 + 6 * dy, 1180, 30, anchor=MIDDLE)
    add_paragraph(tf, [run("pw skill", 10.5, ORANGE, True),
                  run(" enchaîne automatiquement les phases documentaires   "
                      "·   ", 10.5, GRAY), run("pw handoff", 10.5, ORANGE,
                  True), run(" enchaîne les étapes d'implémentation", 10.5,
                  GRAY)], THEME, align=CENTER, first=True, space_after=0)
    footer(s, THEME, "llm-shared - Vue d'ensemble", 9)


def slide_10_doc_phase(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "5. Phase documentaire")
    title(s, THEME, "Phase documentaire : automatisée",
          "pw skill enchaîne les phases sans menu ni intervention humaine "
          "(sauf revue)")
    y0, dy = 158, 66
    skill_tag = [run("== pw skill ==>", 9, ORANGE, True)]
    step_row(s, y0 + 0 * dy, 1, BLUE, LIGHTBLUE, "/write-requirement", BLUE,
             ": écrit le document d'exigence (feature-request ou issue)",
             skill_tag)
    step_row(s, y0 + 1 * dy, 2, ORANGE, WARM, "/review-ask-questions", ORANGE,
             ": l'IA pose des questions ouvertes avec options et "
             "recommandations", [run("[STOP]", 9, WHITE, True)],
             tag_fill=ORANGE)
    step_row(s, y0 + 2 * dy, 3, BLUE, LIGHTBLUE,
             "/consolidate-then-review-ask-questions", BLUE,
             ": plie les réponses dans le document, repose des questions si "
             "besoin", [run("↻ boucle", 9, BLUE, True)])
    step_row(s, y0 + 3 * dy, 4, ORANGE, WARM, "/write-design", ORANGE,
             ": scope, contraintes, comportement cible, scénarios "
             "d'acceptation", skill_tag)
    step_row(s, y0 + 4 * dy, 5, BLUE, LIGHTBLUE, "/write-plans", BLUE,
             ": plan d'exécution numéroté + plan de validation", skill_tag)
    note = rect(s, THEME, 50, y0 + 5 * dy + 4, 1180, 44, fill=LIGHT,
                rounded=True)
    ntf = note.text_frame
    ntf.vertical_anchor = MIDDLE
    add_paragraph(ntf, [run("Le seul arrêt humain : répondre au tableau ", 11,
                  BLUE), run("Q0x | Titre | Réponse recommandée", 11, BLUE,
                  True), run(" lors des revues", 11, BLUE)], THEME,
                  align=CENTER, first=True, space_after=0)
    footer(s, THEME, "llm-shared - Phase documentaire", 10)


def slide_11_impl_phase(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "6. Implémentation")
    title(s, THEME, "Phase d'implémentation : chaîne automatisée",
          "pw handoff enchaîne les étapes sans menu ni « go-ahead » entre les "
          "cases")
    rect(s, THEME, 50, 158, 1180, 210, line=ORANGE, rounded=True)
    tf = textbox(s, THEME, 60, 168, 1160, 22, anchor=MIDDLE)
    add_paragraph(tf, [run("ZONE AUTOMATISÉE : un seul /implement-step N "
                  "déclenche la chaîne", 11, ORANGE, True)], THEME,
                  align=CENTER, first=True, space_after=0)
    box(s, THEME, 90, 205, 500, 66, "/implement-step N", fill=BLUE,
        tcolor=WHITE, size=11, sub="Code + tests · ghog day vert",
        subcolor=TSOFT, subsize=8.5)
    arrow(s, THEME, 595, 225, 90, glyph="▼", size=16)
    box(s, THEME, 690, 205, 500, 66, "/implementation-check N", fill=ORANGE,
        tcolor=WHITE, size=11, sub="Verdict Oui/Non", subcolor=WARM,
        subsize=8.5)
    box(s, THEME, 90, 285, 500, 66, "↻ Non → /implement-missing-step",
        fill=GRAY, tcolor=WHITE, size=10.5,
        sub="Combler les écarts, ghog day, re-check", subcolor=WHITE,
        subsize=8.5)
    arrow(s, THEME, 595, 305, 90, glyph="ou", color=GRAY, size=12)
    box(s, THEME, 690, 285, 500, 66, "✓ Oui → /group-commits-msg", fill=TEAL,
        tcolor=WHITE, size=10.5, sub="Messages de commit groupés → a.commit",
        subcolor=WHITE, subsize=8.5)
    feats = [
        ("🧪 Tests seuils d'abord", "Le plan commence par des tests qui "
         "échouent volontairement. Les étapes suivantes les font passer un "
         "par un.", BLUE),
        ("📦 Commits groupés", "L'IA lit le diff, groupe les fichiers du "
         "moins dépendant au plus dépendant, écrit un message par groupe.",
         ORANGE),
        ("🛑 Portail de commit", "La chaîne s'arrête à a.commit. L'auteur "
         "révise, dit « go ahead », puis gcba rejoue les commits.", BLUE),
    ]
    feature_row(s, THEME, feats, 388)
    footer(s, THEME, "llm-shared - Implémentation", 11)


def slide_12_groundhog(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "7. Groundhog")
    title(s, THEME, "Groundhog : la boucle de tests",
          "Un seul outil remplace les anciens alias : ghog day marche la "
          "chaîne et s'arrête au premier problème")
    box(s, THEME, 540, 158, 200, 40, "ghog day", fill=BLUE, tcolor=WHITE,
        size=13)
    arrow(s, THEME, 630, 200, 20, glyph="▼", size=16)
    steps = [
        ("1. check.bat", "Compilation + lint", LIGHTBLUE, BLUE),
        ("2. Tests affectés", "Rapide, sans couverture", WARM, ORANGE),
        ("3. Suite complète", "Couverture 100% exigée", LIGHTBLUE, BLUE),
    ]
    bw, gap = 220, 40
    total = 3 * bw + 2 * gap
    x0 = (1280 - total) / 2
    for i, (head, sub, fill, line) in enumerate(steps):
        x = x0 + i * (bw + gap)
        st = rect(s, THEME, x, 228, bw, 62, fill=fill, line=line,
                  rounded=True)
        stf = st.text_frame
        stf.vertical_anchor = MIDDLE
        add_paragraph(stf, [run(head, 12, line, True)], THEME, align=CENTER,
                      first=True, space_after=2)
        add_paragraph(stf, [run(sub, 9, MUTED)], THEME, align=CENTER,
                      space_after=0)
        if i < 2:
            arrow(s, THEME, x + bw, 245, gap, glyph="→", size=16)
    arrow(s, THEME, 630, 292, 20, glyph="▼", size=16)
    box(s, THEME, 415, 318, 450, 40, "✓ Objectif atteint : tests verts + "
        "couverture au seuil", fill=TEAL, tcolor=WHITE, size=11)
    feats = [
        ("🔄 Boucle de correction IA", "L'IA lit le code de sortie, applique "
         "le correctif nommé dans le rapport, relance ghog day. Arrêt si "
         "aucun progrès.", BLUE),
        ("📊 Sortie maîtrisée", "Petite pour l'IA (logs redirigés, tail lu). "
         "Vivante pour l'humain (barre de progression, rapport final).",
         ORANGE),
        ("🆔 État de cycle", "Contrat fichier a.ghog.status. L'IA ne devine "
         "jamais : ghog status lit l'état, exit 6 = en cours, 7 = tué.",
         BLUE),
    ]
    feature_row(s, THEME, feats, 392, h=120)
    footer(s, THEME, "llm-shared - Groundhog", 12)


def slide_13_commits(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "8. Commits & Release")
    title(s, THEME, "Commits conventionnels & préparation de release",
          "Au-delà du changelog : le « pourquoi » pour le prochain lecteur")
    banner = rect(s, THEME, 50, 152, 1180, 56, fill=LIGHTBLUE, rounded=True)
    btf = banner.text_frame
    btf.vertical_anchor = MIDDLE
    btf.margin_left = Pt(10)
    btf.margin_right = Pt(10)
    add_paragraph(btf, [run("/group-commits-msg", 11, ORANGE, True),
                  run(" : range les fichiers indexés du moins au plus "
                      "dépendant, écrit un message au format ci-dessous par "
                      "groupe dans a.commit, puis crée un commit par groupe "
                      "après relecture.", 10, BLUE)], THEME, first=True,
                  space_after=0)
    # left column: extended commit format on a dark code panel
    tf = textbox(s, THEME, 50, 224, 570, 22)
    add_paragraph(tf, [run("Format de commit étendu", 12, BLUE, True)], THEME,
                  first=True, space_after=0)
    panel = rect(s, THEME, 50, 252, 570, 240, fill=DARK, rounded=True)
    ptf = panel.text_frame
    ptf.vertical_anchor = TOP
    ptf.margin_left = Pt(12)
    ptf.margin_top = Pt(10)
    add_paragraph(ptf, [run("feat", 10, ORANGE, True),
                  run("(scope): sujet (52 car. max)", 10, CODE)], THEME,
                  first=True, space_after=8)
    add_paragraph(ptf, [run("Why:", 10, GREEN, True)], THEME, space_after=4)
    add_paragraph(ptf, [run("Raison du changement (paragraphe 1).", 10,
                  CODE)], THEME, space_after=1)
    add_paragraph(ptf, [run("État « après » du code (paragraphe 2).", 10,
                  CODE)], THEME, space_after=8)
    add_paragraph(ptf, [run("What:", 10, GREEN, True)], THEME, space_after=4)
    for m in ("- modification 1", "- modification 2", "- modification 3"):
        add_paragraph(ptf, [run(m, 10, CODE)], THEME, space_after=1)
    tf = textbox(s, THEME, 50, 498, 570, 40)
    add_paragraph(tf, [run("Le titre alimente le changelog. Le corps donne le "
                  "contexte que le diff ne porte pas : pour l'humain et pour "
                  "l'IA qui relira plus tard.", 8.5, MUTED)], THEME,
                  first=True, space_after=0, line_spacing=1.1)
    # right column: numbered release steps
    tf = textbox(s, THEME, 660, 224, 570, 22)
    add_paragraph(tf, [run("Préparation de release", 12, BLUE, True)], THEME,
                  first=True, space_after=0)
    steps = [
        ("1", ORANGE, [run("/prepare-release", 10, BLUE, True),
                       run(" : merge, reword, snapshot", 10, TEXT)]),
        ("2", BLUE, [run("/prepare_release_notes", 10, BLUE, True),
                     run(" : titres groupés, résumé, CHANGELOG", 10, TEXT)]),
        ("3", BLUE, [run("Choix du titre de release (3 propositions)", 10,
                         TEXT)]),
        ("4", BLUE, [run("pyproject.toml + uv sync → commit chore(release)",
                         10, TEXT)]),
        ("5", ORANGE, [run("brel", 10, BLUE, True),
                       run(" : build + tag vX.Y.Z [valid]", 10, TEXT)]),
    ]
    y = 256
    for num, col, desc in steps:
        o = oval(s, THEME, 660, y, 30, 30, fill=col)
        otf = o.text_frame
        otf.vertical_anchor = MIDDLE
        add_paragraph(otf, [run(num, 12, WHITE, True)], THEME, align=CENTER,
                      first=True, space_after=0)
        tf = textbox(s, THEME, 700, y, 530, 30, anchor=MIDDLE)
        add_paragraph(tf, desc, THEME, first=True, space_after=0)
        y += 40
    tip = rect(s, THEME, 660, y + 4, 570, 40, fill=WARM, rounded=True)
    ttf = tip.text_frame
    ttf.vertical_anchor = MIDDLE
    add_paragraph(ttf, [run("Le skill ne taggue jamais et ne pousse jamais : "
                  "c'est l'auteur qui décide", 9, ORANGE)], THEME,
                  align=CENTER, first=True, space_after=0)
    footer(s, THEME, "llm-shared - Commits & Release", 13)


def sec_cell(s, x, y, w, h, fill, hcolor, heading, body_runs, accent=False):
    """One security cell from slide 14: tinted rounded box, colored heading,
    body runs, optional orange left accent bar for the 'before prod' gate."""
    rect(s, THEME, x, y, w, h, fill=fill, rounded=True)
    if accent:
        rect(s, THEME, x, y, 5, h, fill=ORANGE)
    tf = textbox(s, THEME, x + 14, y + 8, w - 24, h - 14, anchor=TOP)
    add_paragraph(tf, [run(heading, 11, hcolor, True)], THEME, first=True,
                  space_after=3)
    add_paragraph(tf, body_runs, THEME, space_after=0, line_spacing=1.12)


def slide_14_security(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "Sécurité")
    title(s, THEME, "🛡️ Deux garde-fous sécurité",
          "L'IA accélère le code : elle ne dispense pas des contrôles")
    # two aligned columns; left x=50 w=575, right x=655 w=575, 30px gap
    lx, rx, cw = 50, 655, 575
    tf = textbox(s, THEME, lx, 158, cw, 26, anchor=MIDDLE)
    add_paragraph(tf, [run("🔍 Code et dépendances", 13, BLUE, True)], THEME,
                  align=CENTER, first=True, space_after=0)
    tf = textbox(s, THEME, rx, 158, cw, 26, anchor=MIDDLE)
    add_paragraph(tf, [run("🔒 Fuites de données", 13, ORANGE, True)], THEME,
                  align=CENTER, first=True, space_after=0)
    t = TEXT
    # row 1 : Risque
    sec_cell(s, lx, 192, cw, 84, LIGHTBLUE, BLUE, "⚠️ Risque",
             [run("L'IA peut produire du code vulnérable et tirer des "
                  "bibliothèques anciennes ou trouées.", 10.5, t)])
    sec_cell(s, rx, 192, cw, 84, WARM, ORANGE, "⚠️ Risque",
             [run("Un secret ou une donnée métier collé dans un prompt part "
                  "chez le fournisseur du modèle.", 10.5, t)])
    # row 2 : Contrôle
    sec_cell(s, lx, 288, cw, 84, LIGHTBLUE, BLUE, "✅ Contrôle",
             [run("Passer chaque livraison au scanner ", 10.5, t),
              run("SonarQube ou équivalent", 10.5, t, True),
              run(" : qualité, failles connues, dépendances.", 10.5, t)])
    sec_cell(s, rx, 288, cw, 84, WARM, ORANGE, "🚫 Contrôle",
             [run("Jamais de secrets ni de données réelles dans les prompts",
                  10.5, t, True),
              run(" : clés, tokens, données client, extraits de prod.", 10.5,
                  t)])
    # row 3 : Porte avant prod / Détail
    sec_cell(s, lx, 384, cw, 100, WARM, ORANGE, "🚦 Porte avant prod",
             [run("Avant la mise en production : passage SonarQube + "
                  "correction des vulnérabilités et des bibliothèques. Pas de "
                  "« go » tant que ce n'est pas propre.", 10.5, t)],
             accent=True)
    sec_cell(s, rx, 384, cw, 100, WARM, ORANGE, "🧪 Détail",
             [run("On décrit le problème et la structure, on ne copie pas la "
                  "donnée ; jeux d'exemple anonymisés si besoin.", 10.5, t)])
    # row 4 : Dans le workflow
    sec_cell(s, lx, 496, cw, 84, LIGHTBLUE, BLUE, "🔗 Dans le workflow",
             [run("Après ", 10.5, t), code_run("/implementation-check"),
              run(" et la boucle ", 10.5, t), code_run("groundhog"),
              run(", avant le commit ; le vert SonarQube devient une "
                  "condition de merge.", 10.5, t)])
    sec_cell(s, rx, 496, cw, 84, WARM, ORANGE, "🔗 Dans le workflow",
             [run("Règle qui vaut à toutes les phases : brouillon, exigence, "
                  "design, implémentation.", 10.5, t)])
    # bottom synthesis line
    tf = textbox(s, THEME, 50, 590, 1180, 40, anchor=MIDDLE)
    add_paragraph(tf, [run("💡 L'IA génère vite : ", 10.5, GRAY),
                  run("SonarQube", 10.5, BLUE, True),
                  run(" filtre le code avant prod, la règle ", 10.5, GRAY),
                  run("« jamais de secrets ni de données dans le prompt »",
                      10.5, ORANGE, True),
                  run(" filtre ce qui sort.", 10.5, GRAY)], THEME,
                  align=CENTER, first=True, space_after=0)
    footer(s, THEME, "llm-shared - Sécurité", 14)


def slide_15_benefits(prs):
    s = blank_slide(prs)
    header(s, THEME, BRAND, BRAND_SUB, "Bénéfices")
    title(s, THEME, "Bénéfices pour la DSI",
          "Ce que le framework apporte aux équipes et à la qualité")
    feats = [
        ("📋 Traçabilité complète", "Chaque décision est documentée. Le "
         "« pourquoi » est explicite, pas à déduire du code. Revu par l'IA "
         "avant validation.", BLUE),
        ("🤖 Multi-agent", "GitHub Copilot, Claude Code, ChatGPT Codex. Les "
         "mêmes instructions markdown, les mêmes artefacts.", ORANGE),
        ("🧪 Qualité par construction", "Tests seuils avant implémentation. "
         "Couverture 100% exigée. Revue d'implémentation systématique par "
         "l'IA.", ORANGE),
        ("📚 Documentation vivante", "Brouillon, exigence, design, plan, "
         "validation, commits : chaque phase laisse un artefact relu et "
         "consolidé.", BLUE),
        ("⚡ Automatisation", "pw skill et pw handoff enchaînent les étapes "
         "sans intervention. L'humain intervient seulement aux points de "
         "revue.", BLUE),
        ("🔓 Licence MIT", "Réutilisable dans tous les projets. Copie, fork, "
         "adaptation sans contrainte de copyleft.", ORANGE),
    ]
    w, h, gx, gy = 570, 120, 40, 20
    x0, y0 = 50, 165
    for i, (head, body, accent) in enumerate(feats):
        x = x0 + (i % 2) * (w + gx)
        y = y0 + (i // 2) * (h + gy)
        feature_item(s, THEME, x, y, w, h, head, body, accent=accent)
    footer(s, THEME, "llm-shared - Bénéfices", 15)


def slide_16_conclusion(prs):
    s = blank_slide(prs)
    fill_background(s, BLUE)
    rect(s, THEME, 768, 0, 512, 5, fill=ORANGE)
    brand_top_left(s)
    tf = textbox(s, THEME, 0, 200, 1280, 60, anchor=MIDDLE)
    add_paragraph(tf, [run("Merci de votre attention", 28, WHITE, True)],
                  THEME, align=CENTER, first=True, space_after=0)
    tf = textbox(s, THEME, 240, 285, 800, 70, anchor=MIDDLE)
    add_paragraph(tf, [run("llm-shared - du brouillon à la release", 14,
                  TSOFT)], THEME, align=CENTER, first=True, space_after=4)
    add_paragraph(tf, [run("Sans vibe-coding, avec auto-revue IA et "
                  "traçabilité complète", 11, TMUTED)], THEME, align=CENTER,
                  space_after=0)
    stats = [("6", "phases structurées"), ("3", "agents IA compatibles"),
             ("100%", "traçabilité des décisions")]
    sw, gap = 260, 40
    total = 3 * sw + 2 * gap
    x0 = (1280 - total) / 2
    for i, (num, label) in enumerate(stats):
        x = x0 + i * (sw + gap)
        tf = textbox(s, THEME, x, 400, sw, 50, anchor=MIDDLE)
        add_paragraph(tf, [run(num, 28, ORANGE, True)], THEME, align=CENTER,
                      first=True, space_after=2)
        tf = textbox(s, THEME, x, 452, sw, 24, anchor=MIDDLE)
        add_paragraph(tf, [run(label, 10, TMUTED)], THEME, align=CENTER,
                      first=True, space_after=0)


def main():
    prs = make_presentation(THEME)
    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_problem(prs)
    slide_04_solution(prs)
    slide_05_review(prs)
    slide_06_review_practice(prs)
    slide_07_doc_example(prs)
    slide_08_impl_example(prs)
    slide_09_overview(prs)
    slide_10_doc_phase(prs)
    slide_11_impl_phase(prs)
    slide_12_groundhog(prs)
    slide_13_commits(prs)
    slide_14_security(prs)
    slide_15_benefits(prs)
    slide_16_conclusion(prs)
    prs.save(OUT)
    print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides, "
          f"{os.path.getsize(OUT) // 1024} KB")


if __name__ == "__main__":
    main()


# eof
