"""Deck-specific compositions of the llm-shared deck (copy and adapt per deck).

Each helper draws one repeated visual pattern of the deck -- a phase card, a
flow column, a workflow row, an example panel -- on a SlideCanvas, taking its
geometry as an Area and its variable looks as a small spec dataclass. Keep
generic, deck-agnostic drawing in pptx_helpers.py; keep the patterns tied to
this deck's design language here.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import TYPE_CHECKING

from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from tools.html_to_pptx.deck_theme import (
    BLUE,
    BRAND,
    BRAND_SUB,
    GRAY,
    LIGHTBLUE,
    MUTED,
    ORANGE,
    TMUTED,
    TSOFT,
    WARM,
    WHITE,
)
from tools.html_to_pptx.pptx_helpers import Area, BoxStyle, FrameWriter, run

if TYPE_CHECKING:
    from collections.abc import Sequence

    from pptx.dml.color import RGBColor

    from tools.html_to_pptx.pptx_helpers import Run, SlideCanvas

CENTER = PP_ALIGN.CENTER
MIDDLE = MSO_ANCHOR.MIDDLE


def code_run(text: str) -> Run:
    """An orange slash-command run, matching the HTML inline <code> spans."""
    return run(text, 10, ORANGE)


def phase_box(c: SlideCanvas, area: Area, spec: tuple[str, str, str, str]) -> None:
    """One phase card from slide 4: (num, label, detail, style) on `area`.

    The style is one of "active" (filled blue), "active-orange" (filled
    orange) or "plain" (white with a blue border).
    """
    num, label, detail, style = spec
    if style == "active":
        fill, line, tc, dc = BLUE, BLUE, WHITE, TSOFT
    elif style == "active-orange":
        fill, line, tc, dc = ORANGE, ORANGE, WHITE, WARM
    else:
        fill, line, tc, dc = WHITE, BLUE, BLUE, MUTED
    c.rect(area, fill=fill, line=line, rounded=True)
    writer = c.frame(Area(area.x + 4, area.y + 12, area.w - 8, area.h - 20))
    num_color = ORANGE if style == "plain" else WHITE
    writer.add([run(str(num), 24, num_color, bold=True)], align=CENTER,
               space_after=2)
    writer.add([run(label, 12, tc, bold=True)], align=CENTER, space_after=3)
    writer.add([run(detail, 8.5, dc)], align=CENTER, space_after=0,
               line_spacing=1.05)


@dataclass
class LoopStep:
    """One box of a loop column: its text, fill style and optional note."""

    text: str
    style: str
    note: str | None = None
    note_color: RGBColor | None = None


@dataclass
class LoopLayout:
    """Geometry of a loop column: center x, first box y and box size."""

    cx: float
    y0: float
    box_w: float = 240
    box_h: float = 42


def _loop_step_colors(style: str) -> tuple[RGBColor, RGBColor, RGBColor]:
    """Fill, border and text colors of one loop box style keyword."""
    if style == "filled":
        return BLUE, BLUE, WHITE
    if style == "orange":
        return WHITE, ORANGE, ORANGE
    if style == "filled-orange":
        return ORANGE, ORANGE, WHITE
    return WHITE, BLUE, BLUE


def loop_col(c: SlideCanvas, layout: LoopLayout, title_txt: str,
             sub_txt: Sequence[Run] | None, steps: Sequence[LoopStep]) -> None:
    """A vertical flow column of loop boxes with ▼ connectors and notes."""
    writer = c.frame(Area(layout.cx - 220, layout.y0 - 60, 440, 24),
                     anchor=MIDDLE)
    writer.add([run(title_txt, 13, BLUE, bold=True)], align=CENTER,
               space_after=0)
    if sub_txt:
        sub_writer = c.frame(Area(layout.cx - 240, layout.y0 - 36, 480, 22),
                             anchor=MIDDLE)
        sub_writer.add(sub_txt, align=CENTER, space_after=0)
    y = layout.y0
    for i, st in enumerate(steps):
        fill, line, tc = _loop_step_colors(st.style)
        c.box(Area(layout.cx - layout.box_w / 2, y, layout.box_w,
                   layout.box_h),
              st.text, BoxStyle(fill=fill, line=line, tcolor=tc, size=10))
        y += layout.box_h
        if st.note:
            note_writer = c.frame(
                Area(layout.cx - layout.box_w / 2, y + 1, layout.box_w, 16),
                anchor=MIDDLE)
            note_color = st.note_color if st.note_color is not None else ORANGE
            note_writer.add([run(st.note, 8.5, note_color, bold=True)],
                            align=CENTER, space_after=0)
            y += 17
        if i < len(steps) - 1:
            c.arrow(Area(layout.cx - 10, y, 20, 20))
            y += 20


def wf_row(c: SlideCanvas, y: float, label: str, label_fill: RGBColor,
           desc_runs: Sequence[Run]) -> None:
    """One workflow row from slide 9: colored label + arrow + tinted body."""
    lab = c.rect(Area(50, y, 120, 40), fill=label_fill, rounded=True)
    ltf = lab.text_frame
    ltf.vertical_anchor = MIDDLE
    ltf.word_wrap = True
    lab_writer = FrameWriter(ltf, c.theme)
    lab_writer.add([run(label, 11, WHITE, bold=True)], align=CENTER,
                   space_after=0)
    c.arrow(Area(172, y + 10, 24, 20), glyph="→", size=14)
    bg = LIGHTBLUE if label_fill == BLUE else WARM
    body = c.rect(Area(200, y, 1030, 40), fill=bg, rounded=True)
    btf = body.text_frame
    btf.vertical_anchor = MIDDLE
    btf.word_wrap = True
    btf.margin_left = Pt(8)
    body_writer = FrameWriter(btf, c.theme)
    body_writer.add(desc_runs, space_after=0)


@dataclass
class StepRow:
    """One numbered step row of slide 10: badge, command box and side tag."""

    num: int
    circle_color: RGBColor
    box_fill: RGBColor
    cmd: str
    cmd_color: RGBColor
    desc: str
    tag_runs: Sequence[Run]
    tag_fill: RGBColor | None = None


def step_row(c: SlideCanvas, y: float, spec: StepRow) -> None:
    """Draw one numbered step row: oval badge + command box + tag."""
    o = c.oval(Area(50, y, 44, 44), spec.circle_color)
    otf = o.text_frame
    otf.vertical_anchor = MIDDLE
    o_writer = FrameWriter(otf, c.theme)
    o_writer.add([run(str(spec.num), 16, WHITE, bold=True)], align=CENTER,
                 space_after=0)
    b = c.rect(Area(110, y, 940, 44), fill=spec.box_fill, rounded=True)
    btf = b.text_frame
    btf.vertical_anchor = MIDDLE
    btf.word_wrap = True
    btf.margin_left = Pt(8)
    b_writer = FrameWriter(btf, c.theme)
    b_writer.add([run(spec.cmd, 12, spec.cmd_color, bold=True),
                  run("  " + spec.desc, 10, MUTED)], space_after=0)
    if spec.tag_fill is not None:
        t = c.rect(Area(1075, y + 6, 155, 32), fill=spec.tag_fill,
                   rounded=True)
        ttf = t.text_frame
        ttf.vertical_anchor = MIDDLE
        t_writer = FrameWriter(ttf, c.theme)
        t_writer.add(spec.tag_runs, align=CENTER, space_after=0)
    else:
        tag_writer = c.frame(Area(1075, y, 155, 44), anchor=MIDDLE)
        tag_writer.add(spec.tag_runs, align=CENTER, space_after=0)


def brand_top_left(c: SlideCanvas) -> None:
    """White brand line for the two dark title slides (1 and 16)."""
    writer = c.frame(Area(40, 30, 500, 44))
    writer.add([run(BRAND, 13, WHITE, bold=True)], space_after=1)
    writer.add([run(BRAND_SUB, 7, TMUTED)], space_after=0)


@dataclass
class ExCellStyle:
    """Looks of an example panel: fill, heading color, border and accent."""

    fill: RGBColor
    hcolor: RGBColor
    line: RGBColor | None = None
    accent: RGBColor | None = None
    dashed: bool = False
    hsize: float = 11.5


def ex_cell(c: SlideCanvas, area: Area, heading: str,
            paras: Sequence[Sequence[Run]], style: ExCellStyle) -> None:
    """One example panel from slides 7-8.

    A tinted rounded box with a colored bold heading, then one body
    paragraph per entry in paras (lists of runs).
    """
    r = c.rect(area, fill=style.fill, line=style.line, rounded=True)
    if style.dashed:
        r.line.dash_style = MSO_LINE.DASH
    if style.accent is not None:
        c.rect(Area(area.x, area.y, 5, area.h), fill=style.accent)
    writer = c.frame(Area(area.x + 12, area.y + 8, area.w - 22, area.h - 14))
    writer.add([run(heading, style.hsize, style.hcolor, bold=True)],
               space_after=3)
    for pr in paras:
        writer.add(pr, space_after=3, line_spacing=1.1)


@dataclass
class BannerStyle:
    """Looks of an AVANT/APRÈS band: fill, border and text color."""

    fill: RGBColor
    line: RGBColor | None
    tcolor: RGBColor


def ex_banner(c: SlideCanvas, area: Area, text: tuple[str, str],
              style: BannerStyle, note_runs: Sequence[Run] | None = None) -> None:
    """Full-width AVANT/APRÈS sentence band with a bold (badge, sentence)."""
    badge, sentence = text
    b = c.rect(area, fill=style.fill, line=style.line, rounded=True)
    tf = b.text_frame
    tf.vertical_anchor = MIDDLE
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    writer = FrameWriter(tf, c.theme)
    runs = [run(badge + "   ", 11, style.tcolor, bold=True),
            run(sentence, 11, style.tcolor, italic=True)]
    if note_runs:
        runs += note_runs
    writer.add(runs, space_after=0, line_spacing=1.1)


def ex_connector(c: SlideCanvas, y: float, runs: Sequence[Run]) -> None:
    """Centered ▼ connector line between the example stages."""
    writer = c.frame(Area(50, y, 1180, 20), anchor=MIDDLE)
    writer.add([run("▼  ", 10, GRAY), *runs], align=CENTER, space_after=0)


@dataclass
class SecCellStyle:
    """Looks of a security cell: fill, heading color, optional accent bar."""

    fill: RGBColor
    hcolor: RGBColor
    accent: bool = False


def sec_cell(c: SlideCanvas, area: Area, heading: str,
             body_runs: Sequence[Run], style: SecCellStyle) -> None:
    """One security cell from slide 14.

    A tinted rounded box with a colored heading and body runs, plus an
    optional orange left accent bar for the 'before prod' gate.
    """
    c.rect(area, fill=style.fill, rounded=True)
    if style.accent:
        c.rect(Area(area.x, area.y, 5, area.h), fill=ORANGE)
    writer = c.frame(Area(area.x + 14, area.y + 8, area.w - 24, area.h - 14),
                     anchor=MSO_ANCHOR.TOP)
    writer.add([run(heading, 11, style.hcolor, bold=True)], space_after=3)
    writer.add(body_runs, space_after=0, line_spacing=1.12)


# eof
